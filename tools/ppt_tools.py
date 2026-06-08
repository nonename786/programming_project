# tools/ppt_tools.py
import os
from pathlib import Path
from typing import Dict
from datetime import datetime

from core.tool_registry import register_tool

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


def _get_ppt_path() -> Path:
    """获取PPT存储目录"""
    workspace = Path(os.getenv("MINI_OPENCLAW_WORKSPACE", "workspace")).resolve()
    ppt_dir = workspace / "presentations"
    ppt_dir.mkdir(parents=True, exist_ok=True)
    return ppt_dir


@register_tool(
    name="generate_ppt",
    description="生成PowerPoint演示文稿，支持创建标题页、内容页、列表页等。",
    parameters={
        "type": "object",
        "properties": {
            "title": {"type": "string", "description": "演示文稿标题"},
            "slides": {
                "type": "string",
                "description": "幻灯片内容，JSON格式: [{\"type\": \"title\", \"title\": \"...\", \"subtitle\": \"...\"}, {\"type\": \"content\", \"title\": \"...\", \"content\": \"...\"}]"
            },
            "theme_color": {
                "type": "string",
                "enum": ["blue", "green", "red", "purple"],
                "description": "主题颜色，默认blue"
            },
        },
        "required": ["title", "slides"],
    },
)
def generate_ppt(
        title: str,
        slides: str,
        theme_color: str = "blue",
) -> Dict:
    """生成PPT"""

    if not PPTX_AVAILABLE:
        return {
            "success": False,
            "content": "",
            "error": "python-pptx库未安装，请运行: pip install python-pptx"
        }

    try:
        import json

        # 解析幻灯片数据
        try:
            slides_data = json.loads(slides)
        except json.JSONDecodeError:
            return {"success": False, "content": "", "error": f"slides格式错误，应为JSON数组"}

        if not isinstance(slides_data, list):
            return {"success": False, "content": "", "error": "slides应为JSON数组"}

        # 颜色映射
        color_map = {
            "blue": RGBColor(0, 102, 204),
            "green": RGBColor(0, 153, 76),
            "red": RGBColor(204, 0, 0),
            "purple": RGBColor(102, 0, 153),
        }
        theme_rgb = color_map.get(theme_color, color_map["blue"])

        # 创建演示文稿
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # 添加标题页
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]

        title_shape.text = title
        subtitle_shape.text = f"生成于 {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"

        # 设置标题颜色
        title_shape.text_frame.paragraphs[0].font.color.rgb = theme_rgb
        title_shape.text_frame.paragraphs[0].font.size = Pt(54)

        # 添加内容幻灯片
        for slide_data in slides_data:
            slide_type = slide_data.get("type", "content")

            if slide_type == "title":
                # 标题页
                layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(layout)
                title_shape = slide.shapes.title
                subtitle_shape = slide.placeholders[1]

                title_shape.text = slide_data.get("title", "")
                subtitle_shape.text = slide_data.get("subtitle", "")
                title_shape.text_frame.paragraphs[0].font.color.rgb = theme_rgb

            elif slide_type == "content":
                # 内容页
                layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(layout)
                title_shape = slide.shapes.title
                body_shape = slide.placeholders[1]

                title_shape.text = slide_data.get("title", "")
                title_shape.text_frame.paragraphs[0].font.color.rgb = theme_rgb

                text_frame = body_shape.text_frame
                text_frame.clear()

                content = slide_data.get("content", "")
                if isinstance(content, list):
                    for idx, item in enumerate(content):
                        if idx == 0:
                            p = text_frame.paragraphs[0]
                        else:
                            p = text_frame.add_paragraph()
                        p.text = str(item)
                        p.level = 0
                else:
                    text_frame.text = str(content)

            elif slide_type == "bullet":
                # 列表页
                layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(layout)
                title_shape = slide.shapes.title
                body_shape = slide.placeholders[1]

                title_shape.text = slide_data.get("title", "")
                title_shape.text_frame.paragraphs[0].font.color.rgb = theme_rgb

                text_frame = body_shape.text_frame
                text_frame.clear()

                bullets = slide_data.get("bullets", [])
                for idx, bullet in enumerate(bullets):
                    if idx == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    p.text = str(bullet)
                    p.level = 0

        # 保存PPT
        ppt_dir = _get_ppt_path()
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"presentation_{timestamp}.pptx"
        filepath = ppt_dir / filename
        prs.save(filepath)

        return {
            "success": True,
            "content": f"✅ PPT已生成: {filename}\n📁 路径: {filepath}\n📊 包含 {len(slides_data) + 1} 张幻灯片",
            "error": None
        }

    except Exception as e:
        return {"success": False, "content": "", "error": f"生成PPT失败: {str(e)}"}